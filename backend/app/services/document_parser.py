"""文档解析器 — 多格式文档解析

支持: PDF / DOCX / PPTX / MD / TXT
输出: 分页文本列表 [(page_num, text, metadata), ...]
"""

import hashlib
from pathlib import Path
from dataclasses import dataclass, field


@dataclass
class PageContent:
    page_number: int
    text: str
    metadata: dict = field(default_factory=dict)


@dataclass
class ParsedDocument:
    file_path: str
    file_type: str
    total_pages: int
    pages: list[PageContent]
    file_hash: str
    metadata: dict = field(default_factory=dict)


def _compute_file_hash(file_path: str) -> str:
    sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            sha256.update(chunk)
    return sha256.hexdigest()


class DocumentParser:
    """统一文档解析入口"""

    SUPPORTED_TYPES = {".pdf", ".docx", ".pptx", ".md", ".txt"}

    def parse(self, file_path: str) -> ParsedDocument:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        suffix = path.suffix.lower()
        if suffix not in self.SUPPORTED_TYPES:
            raise ValueError(f"不支持的文件类型: {suffix}")

        parser_map = {
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".pptx": self._parse_pptx,
            ".md": self._parse_text,
            ".txt": self._parse_text,
        }

        pages = parser_map[suffix](file_path)
        file_hash = _compute_file_hash(file_path)

        return ParsedDocument(
            file_path=file_path,
            file_type=suffix,
            total_pages=len(pages),
            pages=pages,
            file_hash=file_hash,
            metadata={"filename": path.name, "size_bytes": path.stat().st_size},
        )

    def _parse_pdf(self, file_path: str) -> list[PageContent]:
        import fitz  # PyMuPDF

        pages = []
        doc = fitz.open(file_path)
        for i, page in enumerate(doc):
            text = page.get_text("text")
            if text.strip():
                pages.append(
                    PageContent(
                        page_number=i + 1,
                        text=text.strip(),
                        metadata={"source": "pdf", "page": i + 1},
                    )
                )
        doc.close()
        return pages

    def _parse_docx(self, file_path: str) -> list[PageContent]:
        from docx import Document

        doc = Document(file_path)
        pages = []
        current_text = []
        page_num = 1

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            if para.style.name.startswith("Heading 1") and current_text:
                pages.append(
                    PageContent(
                        page_number=page_num,
                        text="\n".join(current_text),
                        metadata={"source": "docx"},
                    )
                )
                current_text = []
                page_num += 1

            current_text.append(text)

        if current_text:
            pages.append(
                PageContent(
                    page_number=page_num,
                    text="\n".join(current_text),
                    metadata={"source": "docx"},
                )
            )

        return pages

    def _parse_pptx(self, file_path: str) -> list[PageContent]:
        from pptx import Presentation

        prs = Presentation(file_path)
        pages = []

        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())

            if texts:
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()

                pages.append(
                    PageContent(
                        page_number=i + 1,
                        text="\n".join(texts),
                        metadata={"source": "pptx", "slide_title": title},
                    )
                )

        return pages

    def _parse_text(self, file_path: str) -> list[PageContent]:
        import chardet

        path = Path(file_path)

        with open(file_path, "rb") as f:
            raw = f.read()
            detected = chardet.detect(raw)
            encoding = detected.get("encoding", "utf-8")

        with open(file_path, "r", encoding=encoding, errors="ignore") as f:
            content = f.read()

        if path.suffix.lower() == ".md":
            sections = content.split("\n## ")
            pages = []
            for i, section in enumerate(sections):
                text = section.strip()
                if text:
                    if i > 0:
                        text = "## " + text
                    pages.append(
                        PageContent(
                            page_number=i + 1,
                            text=text,
                            metadata={"source": "markdown"},
                        )
                    )
            return pages if pages else [PageContent(page_number=1, text=content, metadata={"source": "txt"})]
        else:
            lines = content.split("\n")
            page_size = 50
            pages = []
            for i in range(0, len(lines), page_size):
                chunk = "\n".join(lines[i : i + page_size])
                if chunk.strip():
                    pages.append(
                        PageContent(
                            page_number=i // page_size + 1,
                            text=chunk.strip(),
                            metadata={"source": "txt"},
                        )
                    )
            return pages


document_parser = DocumentParser()
